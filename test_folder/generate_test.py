import os
import zipfile
from docx import Document
from docx.shared import Inches
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from openpyxl import Workbook
from openpyxl.drawing.image import Image
from pptx import Presentation
from pptx.util import Inches
from PIL import Image as PILImage
import shutil
import csv

# Tạo hình ảnh giả lập 10MB
def create_large_image(filename):
    img = PILImage.new('RGB', (4000, 3000), color='blue')  # Hình ảnh lớn
    img.save(filename, 'JPEG', quality=95)  # ~10MB
    return filename

# Tạo cấu trúc thư mục
base_dir = "Project_Archive"
dirs = [
    "Data_2025/Documents/Word_Docs/DOCX_Files",
    "Data_2025/Documents/Word_Docs/DOC_Files",
    "Data_2025/Documents/PDF_Storage",
    "Data_2025/Documents/Text_Files",
    "Data_2025/Documents/Presentations",
    "Data_2025/Spreadsheets/XLSX_Records",
    "Data_2025/Spreadsheets/XLS_Archives",
    "Data_2025/Spreadsheets/CSV_Records",
    "Backup/Docs_Backup",
    "Backup/PDF_Backup",
    "Backup/Text_Backup",
    "Backup/Slides_Backup",
    "Miscellaneous/Old_Formats",
    "Miscellaneous/Temp_Files"
]

for d in dirs:
    os.makedirs(os.path.join(base_dir, d), exist_ok=True)

# Tạo hình ảnh 10MB
image_path = "temp_image.jpg"
create_large_image(image_path)

# Hàm tạo file .docx
def create_docx(filename, contains_hello):
    doc = Document()
    doc.add_paragraph("This is a document containing hello" if contains_hello else "This is a sample document")
    doc.add_picture(image_path, width=Inches(5))  # Thêm hình ảnh 10MB
    doc.save(filename)

# Hàm tạo file .doc (giả lập bằng .docx)
def create_doc(filename, contains_hello):
    doc = Document()
    doc.add_paragraph("This is a document containing hello" if contains_hello else "This is a sample document")
    doc.add_picture(image_path, width=Inches(5))
    doc.save(filename)

# Hàm tạo file .pdf
def create_pdf(filename, contains_hello):
    c = canvas.Canvas(filename, pagesize=letter)
    c.drawString(100, 750, "This is a document containing hello" if contains_hello else "This is a sample document")
    c.drawImage(image_path, 100, 500, width=400, height=300)  # Thêm hình ảnh
    c.save()

# Hàm tạo file .xlsx
def create_xlsx(filename, contains_hello):
    wb = Workbook()
    ws = wb.active
    ws['A1'] = "hello" if contains_hello else "sample"
    img = Image(image_path)
    ws.add_image(img, 'B2')
    wb.save(filename)

# Hàm tạo file .xls (giả lập bằng .xlsx)
def create_xls(filename, contains_hello):
    wb = Workbook()
    ws = wb.active
    ws['A1'] = "hello" if contains_hello else "sample"
    img = Image(image_path)
    ws.add_image(img, 'B2')
    wb.save(filename)

# Hàm tạo file .pptx
def create_pptx(filename, contains_hello):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "hello" if contains_hello else "sample"
    slide.shapes.add_picture(image_path, Inches(1), Inches(1), width=Inches(5))
    prs.save(filename)

# Hàm tạo file .txt
def create_txt(filename, contains_hello):
    with open(filename, 'w') as f:
        f.write("hello\n" if contains_hello else "sample\n")
        # Thêm văn bản dài để đạt 10MB
        f.write('A' * (10 * 1024 * 1024))  # ~10MB

# Hàm tạo file .csv
def create_csv(filename, contains_hello):
    with open(filename, 'w', newline='') as f:
        writer = csv.writer(f)
        writer.writerow(["hello"] if contains_hello else ["sample"])
        # Thêm dữ liệu để đạt 10MB
        for _ in range(500000):  # ~10MB
            writer.writerow(['A' * 20])

# Tạo file
file_types = [
    ("Data_2025/Documents/Word_Docs/DOCX_Files", "docx", create_docx),
    ("Data_2025/Documents/Word_Docs/DOC_Files", "doc", create_doc),
    ("Data_2025/Documents/PDF_Storage", "pdf", create_pdf),
    ("Data_2025/Documents/Presentations", "pptx", create_pptx),
    ("Data_2025/Spreadsheets/XLSX_Records", "xlsx", create_xlsx),
    ("Data_2025/Spreadsheets/XLS_Archives", "xls", create_xls),
    ("Data_2025/Documents/Text_Files", "txt", create_txt),
    ("Data_2025/Spreadsheets/CSV_Records", "csv", create_csv)
]

for folder, ext, create_func in file_types:
    for i in range(1, 11):
        filename = os.path.join(base_dir, folder, f"file_{i}.{ext}")
        contains_hello = i <= 7  # 7 file đầu chứa "hello"
        create_func(filename, contains_hello)

# Nén thư mục thành file zip
zip_filename = "Project_Archive.zip"
with zipfile.ZipFile(zip_filename, 'w', zipfile.ZIP_DEFLATED) as zipf:
    for root, _, files in os.walk(base_dir):
        for file in files:
            file_path = os.path.join(root, file)
            arcname = os.path.relpath(file_path, base_dir)
            zipf.write(file_path, os.path.join("Project_Archive", arcname))

# Xóa hình ảnh tạm
os.remove(image_path)

print(f"Đã tạo file zip: {zip_filename}")